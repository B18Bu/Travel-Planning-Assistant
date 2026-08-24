# -*- coding: utf-8 -*-
"""生成《智能文旅策划助手》毕业设计答辩 PPT。

依赖：pip install python-pptx
用法：python build_pptx.py
输出：智能文旅策划助手-答辩.pptx（与脚本同目录）
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = Path(__file__).resolve().parent

# 主题色（与 HTML 版一致）
TEAL = RGBColor(0x1E, 0x6F, 0x66)
TEAL_DARK = RGBColor(0x15, 0x52, 0x4B)
AMBER = RGBColor(0xC9, 0x7A, 0x2B)
INK = RGBColor(0x1D, 0x2B, 0x28)
GRAY = RGBColor(0x57, 0x64, 0x5F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Microsoft YaHei"


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def style_run(run, size: int, bold: bool = False, color: RGBColor = INK) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_cover(prs: Presentation, fields: dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = TEAL_DARK

    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "智能文旅策划助手"
    style_run(r, 54, bold=True, color=WHITE)
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "面向企业场景的只读式智能文旅行程规划 Agent"
    style_run(r, 20, color=RGBColor(0xC9, 0xD4, 0xCD))

    sub = slide.shapes.add_textbox(Inches(1.0), Inches(4.6), Inches(11.3), Inches(2.2))
    stf = sub.text_frame
    stf.word_wrap = True
    lines = [
        ("答辩人", fields["name"]),
        ("学号 / 班级", fields["sid"]),
        ("指导教师", fields["tutor"]),
        ("毕业设计答辩 · 讲解时长 10–15 分钟", ""),
    ]
    first = True
    for label, value in lines:
        p = stf.paragraphs[0] if first else stf.add_paragraph()
        first = False
        p.space_after = Pt(8)
        r = p.add_run()
        if value:
            r.text = f"{label}：{value}"
            style_run(r, 18, color=RGBColor(0xC9, 0xD4, 0xCD))
        else:
            r.text = label
            style_run(r, 18, bold=True, color=RGBColor(0xE2, 0xA4, 0x5C))


def add_slide(prs: Presentation, title: str, items: list[tuple[str, int]], notes: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    # 标题
    t = slide.shapes.title
    t.text_frame.word_wrap = True
    t.text = title
    for p in t.text_frame.paragraphs:
        for r in p.runs:
            style_run(r, 32, bold=True, color=TEAL)
    # 正文
    body = slide.placeholders[1].text_frame
    body.word_wrap = True
    first = True
    for text, level in items:
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(6 if level == 0 else 2)
        r = p.add_run()
        r.text = text
        if text.startswith("📊"):
            style_run(r, 18, color=AMBER)
        else:
            style_run(r, 20 if level == 0 else 18, color=INK if level == 0 else GRAY)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def main() -> None:
    prs = new_presentation()

    add_cover(prs, {"name": "＿＿＿＿＿＿", "sid": "＿＿＿＿＿＿", "tutor": "＿＿＿＿＿＿"})

    add_slide(
        prs,
        "01 项目背景",
        [
            ("业务定位：面向企业内部的只读式行程建议服务", 0),
            ("输入出发地/目的地/日期/人数/天数/预算/偏好 → 天气风险 + 逐日路线 + 住宿 + 餐饮 + 待核验事项", 1),
            ("MVP 只给建议与核验提示，不执行预订、支付或交易", 1),
            ("传统做法痛点", 0),
            ("数据来源异构：和风天气 / 高德地图字段与语义各不相同", 1),
            ("字段漂移：自由字典拼错或改名无人察觉", 1),
            ("来源时间混淆：缓存与实时数据不加区分", 1),
            ("上游错误泄露：原始异常透传给浏览器", 1),
            ("降级状态不一致：局部失败时整体语义不统一", 1),
        ],
        notes="开场 1.5 分钟：一句话定位（强类型数据合同 + 多 Agent 编排 + RAG），再讲痛点。",
    )

    add_slide(
        prs,
        "02 问题分析与核心挑战",
        [
            ("数据源异构 —— 需统一为规范化领域模型", 0),
            ("事实与表达分离 —— itinerary 是唯一事实载体，markdown 不能反推事实", 1),
            ("上游脆弱性 —— 需要缓存 / 重试 / 熔断 / 降级组合", 0),
            ("可追溯性 —— request_id / trace_id 与来源时间语义贯穿", 1),
            ("内容可信 —— 模型只做判断与表达，事实由确定性代码负责", 0),
            ("领域边界 —— 严禁价格 / 库存 / 评分 / 预订等交易字段", 1),
        ],
        notes="强调：难点不在「画出地图数据」，而在把异构脆弱数据收敛为强类型、可追溯、可降级的事实。",
    )

    add_slide(
        prs,
        "03 目标与范围（MVP 边界）",
        [
            ("五个设计目标", 0),
            ("强类型数据合同连接所有模块，杜绝字段漂移", 1),
            ("request_id / trace_id 贯穿一次请求，便于问题定位", 1),
            ("区分实时 / 缓存 / 降级数据，要求用户复核关键事实", 1),
            ("同源双栏页面，供应商原始响应不进浏览器", 1),
            ("服务端密钥 + 固定域名 + 超时重试熔断的最小安全边界", 1),
            ("明确不接入", 0),
            ("OTA、预订、支付、库存、实时价格、评分、订单链接（extra=forbid 直接拒绝）", 1),
        ],
        notes="重点讲非交易边界：从模型层禁止未声明字段，杜绝交易信息混入结果。",
    )

    add_slide(
        prs,
        "04 功能模块总览",
        [
            ("出行规划 —— 天气 → 路线 → 住宿 → 餐饮 → 汇总", 0),
            ("文档知识库 RAG —— 上传解析 / 分块 / 向量化 / 混合检索可溯源", 1),
            ("前端工作台 —— 规划表单 / 知识检索 / 文档库 / 数据看板", 0),
            ("平台能力 —— 韧性三件套 / 追踪标识 / 安全响应头 / 通用错误", 1),
            ("核心模块：api → orchestration → agents → services → models", 0),
            ("强类型合同（models/travel.py）承载全部跨模块数据", 1),
        ],
        notes="1.5 分钟过一遍四大功能域，不需展开细节。",
    )

    add_slide(
        prs,
        "05 技术栈选型",
        [
            ("选型原则：事实性 / 确定性职责交给代码库，判断性 / 表达性才交给模型", 0),
            ("FastAPI + Pydantic v2：extra=forbid 封闭字段，model_validator 跨字段约束", 0),
            ("httpx.AsyncClient + respx：异步调用与同源 mock 测试", 1),
            ("ChromaDB + BGE-small-zh-v1.5（本地）：嵌入式向量库，数据不出域", 0),
            ("PyMuPDF / python-docx：本地解析；Qwen-VL：图表 OCR（仅提取文字）", 1),
            ("DeepSeek：表达润色，客户端已封装、接入设计中", 0),
            ("前端：原生 HTML/JS + 本地 vendor（marked + DOMPurify，SHA-256 校验）", 1),
            ("为什么不引入 LangChain？—— 难点在数据合同与编排韧性，轻量编排器更确定、可测试", 0),
        ],
        notes="3 分钟讲完选型理由，重点回答「为什么不用 LLM 做事实」和「为什么不用 Agent 框架」。",
    )

    add_slide(
        prs,
        "06 系统架构",
        [
            ("📊 本页插入【图 1 · 系统总体架构】（Mermaid 源码见《答辩讲解.md》）", 0),
            ("前端层 —— 同源静态页面，只调相对路径", 0),
            ("API 层 —— FastAPI + 安全中间件：/api/travel-plans、/api/documents、/api/knowledge-search", 1),
            ("编排层 —— SequentialTravelOrchestrator：天气→路线→住宿→餐饮→汇总", 0),
            ("基础设施层 —— 缓存 / 重试 / 熔断 + DocumentStore + ChromaDB", 1),
            ("外部服务 —— 和风天气 / 高德地图 / DeepSeek / Qwen-VL（服务端密钥）", 0),
            ("关键决策：强类型合同贯穿 · 事实与表达分离 · 同源渲染 · 固定域名", 1),
        ],
        notes="图 1 的 Mermaid 源码：flowchart TB，前端/API/编排/Agent/基础设施/外部服务六个子图。讲分层职责。",
    )

    add_slide(
        prs,
        "07 强类型数据合同",
        [
            ("TravelPlanRequest：去空白 / 日期≥今天 / 人数1-20 / 天数1-14 / 预算0-200000 / 偏好≤12", 0),
            ("📊 本页插入【图 2 · Agent 状态机】", 0),
            ("success：有完整 data，无缺失字段", 1),
            ("partial：有 data，必须列出 missing_fields", 1),
            ("degraded：有 data，需说明缺失或受控错误", 1),
            ("failed：data 为空，必有缺失字段与受控错误", 1),
            ("来源时间语义：retrieved_at（获取时间）≠ source_updated_at（供应商真实更新时间）", 0),
        ],
        notes="图 2 的 Mermaid 源码：flowchart LR 状态判定。核心是状态机在每个边界被 model_validator 强制校验。",
    )

    add_slide(
        prs,
        "08 多 Agent 编排与任务流程",
        [
            ("📊 本页插入【图 3 · 旅行规划任务实现流程】", 0),
            ("校验 → 生成 request_id=trace_id → 天气 → 路线 → 住宿 → 餐饮 → 汇总 → 前端渲染", 1),
            ("单 Agent 异常由 _safe_agent_call 转为受控 failed，其余 Agent 继续执行", 0),
            ("天气结果（含降级）作为路线 Agent 输入 —— 按天气风险调整行程", 1),
            ("输出 TravelPlanDocument：状态 + 来源 + 待核验 + Markdown", 0),
        ],
        notes="图 3 的 Mermaid 源码：flowchart TB 完整链路。重点演示降级传递路径。",
    )

    add_slide(
        prs,
        "五类专业 Agent 的职责",
        [
            ("WeatherAgent —— 地理编码 + 3 日预报；风险等级 / 出行提醒 / 室内偏好", 0),
            ("RouteAgent —— 往返驾车预估；按天气风险选景区；每日 上午/下午/傍晚 · 120 分钟", 0),
            ("LodgingAgent —— 推荐区域住宿 POI；仅位置与筛选建议，无交易字段", 1),
            ("FoodAgent —— 上午景区→午餐，傍晚/下午→晚餐；附近 2km 餐饮 POI", 0),
            ("SummaryAgent —— 聚合去重 / 定顶层状态 / 确定性模板生成 Markdown", 1),
            ("原则：不复用候选、不伪造结果，候选不足时明确列出待补字段", 0),
        ],
        notes="逐个讲职责与输入输出，强调确定性、可核验。",
    )

    add_slide(
        prs,
        "09 文档知识库 · RAG 混合检索",
        [
            ("📊 本页插入【图 4 · 文档知识库流水线】", 0),
            ("入库：上传校验 → 解析（PyMuPDF/python-docx/Qwen-VL OCR）→ 重叠分块 → BGE 向量化 → ChromaDB", 1),
            ("检索：查询解析（城市 + 意图）→ 语义检索 + 关键词检索 → RRF 融合 → 可溯源命中", 0),
            ("城市硬过滤：避免「语义相近但城市不符」的误命中", 1),
            ("一致性：原文在 DocumentStore，Chroma 只存向量；未 ready 不参与检索", 0),
            ("每个分块保留 page/section/table/figure 来源定位", 1),
        ],
        notes="图 4 的 Mermaid 源码：flowchart TB 入库链路 + 检索链路。重点讲城市硬过滤与双写一致性。",
    )

    add_slide(
        prs,
        "10 韧性设计 · 缓存 / 重试 / 熔断",
        [
            ("📊 本页插入【图 5 · 外部请求韧性流程】", 0),
            ("缓存命中 → data_status=cached；上游成功 → realtime + 记录 source_updated_at", 1),
            ("仅对连接错误 / 超时 / 429 / 5xx 做最多 3 次指数退避重试", 0),
            ("熔断：连续失败达阈值打开 60s，generation token 控制半开探测", 1),
            ("TTL 按数据变化频率：地理编码 7 天 / 路线 15 分钟 / POI 1 小时 / 天气 30 分钟", 0),
            ("如实说明已知限制：3 日预报窗口 / 顺序执行无总 deadline / 生产化治理未实现", 1),
        ],
        notes="图 5 的 Mermaid 源码：flowchart TB。诚实交代局限，能增加可信度。",
    )

    add_slide(
        prs,
        "11 安全边界",
        [
            ("密钥与域名：密钥仅后端读取；base_url 固定校验，客户端无法覆盖", 0),
            ("严格模型：extra=forbid 拒绝未声明字段；frozen 不可变", 1),
            ("XSS 防护：marked + DOMPurify 净化；FORBID_TAGS / FORBID_ATTR 禁脚本与事件属性", 0),
            ("供应链：前端依赖本地 vendor + SHA-256，不用 CDN", 1),
            ("通用错误：不泄露堆栈、参数与供应商原始响应", 0),
            ("上传安全：MIME/后缀/签名校验；文件名安全基名；MinerU 仅 HTTPS 白名单", 1),
        ],
        notes="安全是六条线同时生效，不是单个开关。",
    )

    add_slide(
        prs,
        "12 前端与交互",
        [
            ("首页工作台：规划表单 + 数字人讲解占位 + 知识检索", 0),
            ("任务结果：左行程正文 / 右待核验事项、来源与更新时间、降级说明", 1),
            ("文档库：上传、处理状态、分块浏览、删除管理", 0),
            ("数据看板：赞踩反馈与好评率（当前为界面原型 · 演示数据）", 1),
            ("响应式：768px 断点双栏转单栏", 0),
            ("可信交互：每条建议伴随「待核验 / 来源时间」元信息", 1),
        ],
        notes="强调可信交互原则：JSON 仅 API 消费者与测试使用，前端从不直接暴露供应商响应。",
    )

    add_slide(
        prs,
        "13 测试与质量保障",
        [
            ("respx 拦截 httpx：模拟成功 / 超时 / 429 / 5xx，零真实调用", 0),
            ("合同校验测试：验证「为何需要此行为」而非仅「返回了什么」", 1),
            ("覆盖：模型 / 配置 / 客户端 / 韧性 / 四类 Agent / 汇总 / 编排 / API / 前端资源", 0),
            ("文档管线：分块边界、来源定位、批量上传状态、删除补偿", 1),
            ("质量门禁：文档测试 + 全量测试 + git diff --check", 0),
        ],
        notes="强调测试意图（规则 9）：测试验证为什么需要此行为。",
    )

    add_slide(
        prs,
        "14 项目亮点与创新点",
        [
            ("强类型数据合同驱动 —— 字段漂移 / 状态不一致消灭在类型层", 0),
            ("确定性事实 + 受控表达 —— 从架构上防 LLM 幻觉", 1),
            ("四态状态机 + 来源时间语义 —— 数据诚实可核验", 0),
            ("韧性三件套 + 降级贯穿 —— 单点故障不拖垮整体", 1),
            ("可溯源混合检索 RAG —— 城市硬过滤 + RRF 融合", 0),
            ("领域边界清晰 —— 非交易边界从模型层强制", 1),
        ],
        notes="2 分钟讲亮点，每条一句话 + 一句佐证。",
    )

    add_slide(
        prs,
        "15 当前局限与后续规划",
        [
            ("近期迭代：DeepSeek 表达润色接入；3 日窗口之外支持；扫描件 OCR / MinerU 落地", 0),
            ("性能优化：single-flight 防击穿；连接池复用；缓存淘汰；请求级总超时预算", 1),
            ("生产化：认证 / 限流 / 租户隔离；集中式密钥；审计日志与监控告警", 0),
            ("体验扩展：数字人讲解；知识库反馈闭环；多文档对比问答", 1),
        ],
        notes="承认局限 + 给出路径，比回避更稳妥。",
    )

    add_slide(
        prs,
        "16 讲解节奏（10–15 分钟）",
        [
            ("0:00–1:30 背景与痛点", 0),
            ("1:30–3:00 目标范围 + 功能模块", 1),
            ("3:00–5:00 技术栈选型", 0),
            ("5:00–7:30 架构 + 数据合同 + 状态机", 1),
            ("7:30–10:00 编排流程 + 五类 Agent", 0),
            ("10:00–11:30 RAG + 韧性", 1),
            ("11:30–12:30 安全 + 前端 + 测试", 0),
            ("12:30–14:00 亮点 + 局限", 1),
            ("14:00–15:00 总结致谢 + 提问缓冲", 0),
        ],
        notes="自用时控参考，正式答辩可去掉本页。",
    )

    add_slide(
        prs,
        "17 答辩题库速览（附录）",
        [
            ("系统设计：为什么多 Agent 顺序编排？为什么不用 LLM 做事实？状态机怎么设计？", 0),
            ("技术细节：熔断半开如何实现？TTL 为什么分粒度？哪些错误不重试？", 1),
            ("RAG：为什么混合检索 + 城市硬过滤？Chroma 与 DocumentStore 为何分开？", 0),
            ("安全：XSS 防护链路？密钥为什么不进前端？", 1),
            ("工程：如何零费用测试外部 API？多租户要改哪些？最大难点？", 0),
            ("完整 19 题参考回答见《答辩讲解.md》", 1),
        ],
        notes="备用页，评委问完再展开；答案在 Markdown / HTML 版题库。",
    )

    add_slide(
        prs,
        "18 结语",
        [
            ("事实靠确定性代码保证正确，表达用受控 LLM 增强可读性", 0),
            ("强类型数据合同 + 多 Agent 编排 + RAG 知识库 + 韧性 + 安全边界", 1),
            ("当前是可核验、可测试、安全边界的 MVP", 0),
            ("后续：表达润色 / 生产化治理 / 体验扩展", 1),
            ("以上讲解完毕，欢迎各位老师提问。", 0),
        ],
        notes="收尾致谢，语速放慢。",
    )

    out = HERE / "智能文旅策划助手-答辩.pptx"
    prs.save(out)
    print(f"已生成：{out}")


if __name__ == "__main__":
    main()
